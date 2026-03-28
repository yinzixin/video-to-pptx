"""Build a .pptx from a typed slide plan and frame images.

Each slide_type gets a distinct visual layout with its own accent colour,
title bar, and content arrangement so the resulting deck is varied and
pedagogically clear.
"""

from __future__ import annotations

import os
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from slide_plan import SlideSpec, SlidePlan

# ---------------------------------------------------------------------------
# Colour palette — one accent per slide type
# ---------------------------------------------------------------------------

_C = RGBColor

PALETTE: dict[str, dict[str, RGBColor]] = {
    "story_intro":   {"bar": _C(0x1B, 0x3A, 0x5C), "light": _C(0xE0, 0xE8, 0xF0)},
    "plot_summary":  {"bar": _C(0x2D, 0x6A, 0x4F), "light": _C(0xE2, 0xF0, 0xE8)},
    "key_scene":     {"bar": _C(0xE0, 0x6C, 0x4E), "light": _C(0xFD, 0xEE, 0xE9)},
    "vocabulary":    {"bar": _C(0x6C, 0x3C, 0x97), "light": _C(0xEE, 0xE4, 0xF6)},
    "key_phrases":   {"bar": _C(0x2A, 0x9D, 0x8F), "light": _C(0xE2, 0xF4, 0xF1)},
    "comprehension": {"bar": _C(0xC4, 0x96, 0x1A), "light": _C(0xFA, 0xF3, 0xDE)},
    "moral_lesson":  {"bar": _C(0xB5, 0x1A, 0x2B), "light": _C(0xF6, 0xE1, 0xE3)},
    "discussion":    {"bar": _C(0x45, 0x7B, 0x9D), "light": _C(0xE4, 0xEE, 0xF4)},
}

WHITE = _C(0xFF, 0xFF, 0xFF)
DARK = _C(0x2C, 0x2C, 0x2C)
MID_GRAY = _C(0x66, 0x66, 0x66)
LIGHT_BG = _C(0xF5, 0xF5, 0xF5)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

SECTION_TAGS: dict[str, str] = {
    "story_intro": "Meet the Characters",
    "plot_summary": "What Happens",
    "key_scene": "Key Scene",
    "vocabulary": "Word Bank",
    "key_phrases": "Useful Expressions",
    "comprehension": "Check Understanding",
    "moral_lesson": "The Lesson",
    "discussion": "Let's Discuss",
}

# ---------------------------------------------------------------------------
# Low‑level helpers
# ---------------------------------------------------------------------------


def _colors(slide_type: str) -> dict[str, RGBColor]:
    return PALETTE.get(slide_type, PALETTE["key_scene"])


def _add_colored_bar(
    slide: Any,
    color: RGBColor,
    title: str,
    tag: str = "",
) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(7.5), Inches(0.55))
    tf = tb.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    if tag:
        tb2 = slide.shapes.add_textbox(
            Inches(7.8), Inches(0.15), Inches(2.0), Inches(0.42)
        )
        tf2 = tb2.text_frame
        tf2.text = tag
        tf2.paragraphs[0].font.size = Pt(11)
        tf2.paragraphs[0].font.color.rgb = _C(0xDD, 0xDD, 0xDD)
        tf2.paragraphs[0].font.italic = True
        tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT


def _add_teacher_notes(slide: Any, notes: str | None) -> None:
    if not notes:
        return
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(4.85), SLIDE_W, Inches(0.775)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_BG
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.88), Inches(9.0), Inches(0.7)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = f"Teacher tip: {notes}"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = MID_GRAY


def _add_bullets(
    slide: Any,
    bullets: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: int = 16,
    color: RGBColor = DARK,
    spacing: int = 6,
) -> None:
    if not bullets:
        return
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        if i == 0:
            tf.text = line
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_before = Pt(spacing)


def _add_image_safe(
    slide: Any,
    path: str | None,
    left: float,
    top: float,
    *,
    max_width: float | None = None,
    max_height: float | None = None,
) -> Any | None:
    """Add an image to the slide. Returns the picture shape or None."""
    if not path or not os.path.isfile(path):
        return None
    kwargs: dict[str, Any] = {}
    if max_height is not None:
        kwargs["height"] = Inches(max_height)
    elif max_width is not None:
        kwargs["width"] = Inches(max_width)
    return slide.shapes.add_picture(
        path, Inches(left), Inches(top), **kwargs
    )


def _blank_slide(prs: Presentation) -> Any:
    idx = 6 if len(prs.slide_layouts) > 6 else len(prs.slide_layouts) - 1
    return prs.slides.add_slide(prs.slide_layouts[idx])


# ---------------------------------------------------------------------------
# Type‑specific renderers
# ---------------------------------------------------------------------------


def _render_story_intro(
    prs: Presentation, spec: SlideSpec, image_path: str | None
) -> None:
    slide = _blank_slide(prs)
    c = _colors("story_intro")
    _add_colored_bar(slide, c["bar"], spec.title, SECTION_TAGS["story_intro"])

    body_top = 0.85
    if image_path and os.path.isfile(image_path):
        _add_image_safe(slide, image_path, 0.4, body_top, max_height=3.6)
        _add_bullets(
            slide, spec.bullets, 5.2, body_top, 4.3, 3.6, font_size=15
        )
    else:
        _add_bullets(slide, spec.bullets, 0.5, body_top, 9.0, 3.8, font_size=16)

    _add_teacher_notes(slide, spec.teacher_notes)


def _render_plot_summary(
    prs: Presentation, spec: SlideSpec, image_path: str | None
) -> None:
    slide = _blank_slide(prs)
    c = _colors("plot_summary")
    _add_colored_bar(slide, c["bar"], spec.title, SECTION_TAGS["plot_summary"])

    if image_path and os.path.isfile(image_path):
        _add_image_safe(slide, image_path, 6.8, 0.9, max_height=2.6)
        _add_bullets(
            slide, spec.bullets, 0.5, 0.85, 6.0, 3.8, font_size=15, spacing=8
        )
    else:
        _add_bullets(
            slide, spec.bullets, 0.5, 0.85, 9.0, 3.8, font_size=16, spacing=8
        )

    _add_teacher_notes(slide, spec.teacher_notes)


def _render_key_scene(
    prs: Presentation, spec: SlideSpec, image_path: str | None
) -> None:
    slide = _blank_slide(prs)
    c = _colors("key_scene")
    _add_colored_bar(slide, c["bar"], spec.title, SECTION_TAGS["key_scene"])

    next_top = 0.85

    if image_path and os.path.isfile(image_path):
        pic = _add_image_safe(slide, image_path, 0.5, next_top, max_height=2.2)
        if pic:
            context_left = (pic.left + pic.width) / 914400 + 0.3
            context_width = 9.5 - context_left
        else:
            context_left = 0.5
            context_width = 9.0
    else:
        context_left = 0.5
        context_width = 9.0

    if spec.bullets:
        _add_bullets(
            slide,
            spec.bullets,
            context_left,
            next_top,
            context_width,
            1.0,
            font_size=14,
            color=MID_GRAY,
        )

    dialogue = spec.scene_dialogue or []
    if dialogue:
        dlg_top = 3.3 if (image_path and os.path.isfile(image_path)) else 2.0
        dlg_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(dlg_top), Inches(8.6), Inches(1.8)
        )
        tf = dlg_box.text_frame
        tf.word_wrap = True
        for j, dl in enumerate(dialogue):
            if j == 0:
                tf.text = ""
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run_speaker = p.add_run()
            run_speaker.text = f"{dl.speaker}: "
            run_speaker.font.bold = True
            run_speaker.font.size = Pt(15)
            run_speaker.font.color.rgb = c["bar"]
            run_line = p.add_run()
            run_line.text = dl.line
            run_line.font.size = Pt(15)
            run_line.font.color.rgb = DARK
            p.space_before = Pt(6)

    _add_teacher_notes(slide, spec.teacher_notes)


def _render_vocabulary(
    prs: Presentation, spec: SlideSpec, image_path: str | None
) -> None:
    slide = _blank_slide(prs)
    c = _colors("vocabulary")
    _add_colored_bar(slide, c["bar"], spec.title, SECTION_TAGS["vocabulary"])

    items = spec.vocab_items or []
    if not items:
        _add_bullets(slide, spec.bullets, 0.5, 0.85, 9.0, 3.8, font_size=16)
        _add_teacher_notes(slide, spec.teacher_notes)
        return

    display_items = items[:8]
    rows = len(display_items) + 1
    tbl_shape = slide.shapes.add_table(
        rows, 4, Inches(0.3), Inches(0.85), Inches(9.4), Inches(rows * 0.45)
    )
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(1.6)
    tbl.columns[1].width = Inches(0.9)
    tbl.columns[2].width = Inches(2.8)
    tbl.columns[3].width = Inches(4.1)

    headers = ["Word", "POS", "Definition", "Example"]
    for j, hdr in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = c["bar"]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.LEFT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08)

    for i, vi in enumerate(display_items):
        vals = [vi.word, vi.pos, vi.definition, vi.example]
        for j, val in enumerate(vals):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = c["light"]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK
                if j == 0:
                    p.font.bold = True
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)

    _add_teacher_notes(slide, spec.teacher_notes)


def _render_key_phrases(
    prs: Presentation, spec: SlideSpec, image_path: str | None
) -> None:
    slide = _blank_slide(prs)
    c = _colors("key_phrases")
    _add_colored_bar(slide, c["bar"], spec.title, SECTION_TAGS["key_phrases"])

    _add_bullets(
        slide, spec.bullets, 0.8, 1.0, 8.4, 3.5, font_size=18, spacing=12
    )
    _add_teacher_notes(slide, spec.teacher_notes)


def _render_comprehension(
    prs: Presentation, spec: SlideSpec, image_path: str | None
) -> None:
    slide = _blank_slide(prs)
    c = _colors("comprehension")
    _add_colored_bar(slide, c["bar"], spec.title, SECTION_TAGS["comprehension"])

    if not spec.bullets:
        _add_teacher_notes(slide, spec.teacher_notes)
        return

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9.0), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(spec.bullets):
        is_question = line.strip().upper().startswith("Q:")
        if i == 0:
            tf.text = ""
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.space_before = Pt(8 if is_question else 2)
        if is_question:
            p.font.bold = True
            p.font.color.rgb = c["bar"]
        else:
            p.font.color.rgb = DARK
            p.font.italic = line.strip().upper().startswith("A:")

    _add_teacher_notes(slide, spec.teacher_notes)


def _render_moral_lesson(
    prs: Presentation, spec: SlideSpec, image_path: str | None
) -> None:
    slide = _blank_slide(prs)
    c = _colors("moral_lesson")
    _add_colored_bar(slide, c["bar"], spec.title, SECTION_TAGS["moral_lesson"])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8),
        Inches(1.2),
        Inches(8.4),
        Inches(3.0),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = c["light"]
    bg.line.fill.background()

    if spec.bullets:
        tb = slide.shapes.add_textbox(
            Inches(1.2), Inches(1.5), Inches(7.6), Inches(2.5)
        )
        tf = tb.text_frame
        tf.word_wrap = True

        tf.text = spec.bullets[0]
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = c["bar"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        for line in spec.bullets[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(15)
            p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(10)

    if image_path and os.path.isfile(image_path):
        _add_image_safe(slide, image_path, 7.5, 3.6, max_height=1.5)

    _add_teacher_notes(slide, spec.teacher_notes)


def _render_discussion(
    prs: Presentation, spec: SlideSpec, image_path: str | None
) -> None:
    slide = _blank_slide(prs)
    c = _colors("discussion")
    _add_colored_bar(slide, c["bar"], spec.title, SECTION_TAGS["discussion"])

    _add_bullets(
        slide, spec.bullets, 0.8, 1.0, 8.4, 3.5, font_size=17, spacing=14
    )
    _add_teacher_notes(slide, spec.teacher_notes)


_RENDERERS: dict[str, Any] = {
    "story_intro": _render_story_intro,
    "plot_summary": _render_plot_summary,
    "key_scene": _render_key_scene,
    "vocabulary": _render_vocabulary,
    "key_phrases": _render_key_phrases,
    "comprehension": _render_comprehension,
    "moral_lesson": _render_moral_lesson,
    "discussion": _render_discussion,
}


def _render_generic(
    prs: Presentation, spec: SlideSpec, image_path: str | None
) -> None:
    """Fallback for any unknown slide_type."""
    slide = _blank_slide(prs)
    c = _colors(spec.slide_type)
    _add_colored_bar(slide, c["bar"], spec.title, spec.slide_type.replace("_", " ").title())

    body_top = 0.85
    if image_path and os.path.isfile(image_path):
        _add_image_safe(slide, image_path, 0.5, body_top, max_height=3.2)
        _add_bullets(slide, spec.bullets, 5.3, body_top, 4.2, 3.6, font_size=15)
    else:
        _add_bullets(slide, spec.bullets, 0.5, body_top, 9.0, 3.8, font_size=16)

    _add_teacher_notes(slide, spec.teacher_notes)


# ---------------------------------------------------------------------------
# Title slide and overview
# ---------------------------------------------------------------------------


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _add_overview_slide(
    prs: Presentation,
    plan: SlidePlan,
) -> None:
    """Lesson overview: story summary, moral, learning objectives."""
    slide = _blank_slide(prs)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _C(0x34, 0x49, 0x5E)
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9.0), Inches(0.55))
    tf = tb.text_frame
    tf.text = "Lesson Overview"
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    content_top = 0.85
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(content_top), Inches(9.0), Inches(4.5)
    )
    btf = box.text_frame
    btf.word_wrap = True

    sections: list[tuple[str, list[str]]] = []
    if plan.story_summary:
        sections.append(("Story Summary", [plan.story_summary]))
    if plan.moral:
        sections.append(("The Moral", [plan.moral]))
    if plan.learning_objectives:
        sections.append(
            ("Learning Objectives", [f"  {o}" for o in plan.learning_objectives[:8]])
        )
    if plan.teaching_rationale:
        sections.append(("Teaching Rationale", [plan.teaching_rationale]))

    first = True
    for heading, lines in sections:
        if first:
            btf.text = ""
            p = btf.paragraphs[0]
            first = False
        else:
            p = btf.add_paragraph()
            p.text = ""
            p.space_before = Pt(6)
            p = btf.add_paragraph()

        p.text = heading
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = _C(0x34, 0x49, 0x5E)

        for line in lines:
            p2 = btf.add_paragraph()
            p2.text = line
            p2.font.size = Pt(13)
            p2.font.color.rgb = DARK
            p2.space_before = Pt(2)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------


def build_presentation(
    plan: SlidePlan,
    frames_manifest: dict[str, Any],
    video_basename: str,
    out_path: str,
) -> None:
    frames = frames_manifest.get("frames", [])
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _add_title_slide(prs, plan.lesson_title, f"Source: {video_basename}")

    has_overview = (
        plan.story_summary
        or plan.moral
        or plan.learning_objectives
        or plan.teaching_rationale
    )
    if has_overview:
        _add_overview_slide(prs, plan)

    for spec in plan.slides:
        img_path: str | None = None
        if frames and 0 <= spec.frame_index < len(frames):
            img_path = frames[spec.frame_index].get("path")

        renderer = _RENDERERS.get(spec.slide_type, _render_generic)
        renderer(prs, spec, img_path)

    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)
